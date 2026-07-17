"""
End-to-end integration test for the branded-pptx skill ON SERVERLESS (issue #2).

The counterpart to e2e_test.py, but for a BUILDER skill: instead of a markdown report, the
job must produce a real .pptx. This is the test that #2's acceptance turns on - the skill was
once merged on local-only proof, so "it builds a deck on my laptop" is explicitly not enough:

    put a markdown document into the INPUT volume
      -> trigger the deployed job with --skill-dir <volume>/skills/branded-pptx
        -> wait for a terminal state
          -> download the .pptx from the OUTPUT volume and REOPEN it with python-pptx
            -> assert real slides, then clean up

Reopening the artifact is the point: a file of the right name proves nothing (a zero-byte or
corrupt file would pass an existence check). If python-pptx can parse it, PowerPoint can too.

Exit code 0 = PASS, 1 = FAIL. Run after `databricks bundle deploy`:

    python scripts/e2e_pptx_test.py --profile coldstart
"""
import argparse
import datetime
import io
import json
import os
import subprocess
import sys
import time

from databricks.sdk import WorkspaceClient
from pptx import Presentation

# A benign document with the structure the skill maps to slides: H1 -> title slide, each H2 ->
# a content slide, list items -> bullets. Deliberately free of anything the content guardrail
# would flag, so this test exercises the builder path and not the reject queue.
#
# Deliberately NO run token in the text: the LLM guard reads the DOCUMENT only, and a numeric
# token like 20260716-203317 can read as an identifier, quarantining the input so the job
# returns SUCCESS with no deck - a flaky suite that misdiagnoses as a builder bug. The token
# lives in the filename (which the guard never sees), and that is what keeps runs distinct.
DOC = """# Platform Update

Serverless skills now build decks without a design toolchain.

## Highlights

- Skills publish once to a shared volume and any job consumes them
- The runner supplies the plumbing and the skill supplies the behavior

## Next

- Chain a report into a deck in a single job
"""


def resolve_job_id(profile: str, resource_key: str) -> int:
    out = subprocess.run(
        ["databricks", "bundle", "summary", "-o", "json", "-p", profile],
        capture_output=True, text=True, check=True,
    ).stdout
    jobs = json.loads(out).get("resources", {}).get("jobs", {})
    if resource_key not in jobs:
        raise SystemExit(f"FAIL: job resource '{resource_key}' not found in bundle summary. "
                         f"Did you run `databricks bundle deploy -p {profile}` from this folder?")
    return int(jobs[resource_key]["id"])


def override(params, flag, value):
    """Set a flag on the DEPLOYED parameter list, preserving every other deployed param.

    A wholesale python_params replacement would silently drop --rejected-dir and friends; the
    job's own deployed values are the baseline and we redirect only what this test owns.
    """
    p = list(params)
    if flag in p:
        p[p.index(flag) + 1] = value
    else:
        p += [flag, value]
    return p


def main():
    ap = argparse.ArgumentParser(description="Black-box e2e test of branded-pptx on serverless.")
    ap.add_argument("--profile", default=os.environ.get("DATABRICKS_CONFIG_PROFILE", "coldstart"))
    ap.add_argument("--catalog", default=os.environ.get("DATABRICKS_CATALOG", "workspace"))
    ap.add_argument("--schema", default=os.environ.get("DATABRICKS_SCHEMA", "genai"))
    ap.add_argument("--resource-key", default="mvp0_weekly_report")
    ap.add_argument("--timeout-min", type=int, default=10)
    ap.add_argument("--keep", action="store_true", help="Do not delete the test files afterward.")
    args = ap.parse_args()

    w = WorkspaceClient(profile=args.profile)
    job_id = resolve_job_id(args.profile, args.resource_key)

    base = f"/Volumes/{args.catalog}/{args.schema}"
    token = datetime.datetime.now().strftime("%Y%m%d-%H%M%S")
    stem = f"pptx-e2e-{token}"
    in_path = f"{base}/input/{stem}.md"
    out_dir = f"{base}/output"
    skill_dir = f"{base}/skills/branded-pptx"
    started = time.time()

    def step(msg):
        print(f"[pptx +{time.time()-started:5.1f}s] {msg}")

    step(f"job_id={job_id}  skill={skill_dir}")

    # 1) PUT a real markdown document into the input volume.
    w.files.upload(in_path, io.BytesIO(DOC.encode("utf-8")), overwrite=True)
    step(f"uploaded {in_path}")

    # Everything after the upload runs under try/finally so the shared volumes are cleaned up on
    # EVERY exit - including the one that matters most: a corrupt deck that fails to reopen (the
    # exact regression this test exists to catch) raises out of Presentation(), and without this
    # the debris would pile up on the volume precisely when the test is doing its job.
    out_path = None
    try:
        # 2) TRIGGER the deployed job pointed at the branded-pptx skill on the volume.
        deployed = list(w.jobs.get(job_id=job_id).settings.tasks[0].spark_python_task.parameters or [])
        params = override(deployed, "--skill-dir", skill_dir)
        params = override(params, "--in-path", in_path)
        params = override(params, "--out-dir", out_dir)
        step("triggering job run and waiting for terminal state...")
        run = w.jobs.run_now(job_id=job_id, python_params=params).result(
            timeout=datetime.timedelta(minutes=args.timeout_min))
        state = run.state.result_state.value if run.state and run.state.result_state else "UNKNOWN"
        step(f"run finished: result_state={state}  url={run.run_page_url}")
        if state != "SUCCESS":
            print(f"\nRESULT: FAIL - job run did not succeed (state={state})")
            return 1

        # 3) GET the deck. The skill names it <stem>-branded-pptx-<date>.pptx (skill-namespaced).
        expected = f"{stem}-branded-pptx"
        matches = [e.path for e in w.files.list_directory_contents(out_dir)
                   if e.path.rsplit("/", 1)[-1].startswith(expected) and e.path.endswith(".pptx")]
        if len(matches) != 1:
            # Name the likeliest cause instead of leaving a bare "no deck". The content guard can
            # quarantine the input and return early, so the job SUCCEEDS with the skill never run -
            # a real asymmetry that would otherwise read as a builder bug.
            quarantined = [e.path for e in w.files.list_directory_contents(f"{base}/rejected")
                           if e.path.rsplit("/", 1)[-1].startswith(stem)]
            if quarantined:
                print(f"\nRESULT: FAIL - the content guard quarantined the test input ({quarantined}); "
                      f"the skill never ran. Make the test document less identifier-like.")
            else:
                print(f"\nRESULT: FAIL - expected exactly one '{expected}*.pptx' in the output, got {matches}")
            return 1
        out_path = matches[0]
        data = w.files.download(out_path).contents.read()
        step(f"downloaded {out_path} ({len(data)} bytes)")

        # 4) ASSERT it is a REAL deck by reopening it - the whole point of this test. A name check
        #    would pass on a truncated file; python-pptx parsing it means PowerPoint can open it.
        prs = Presentation(io.BytesIO(data))
        text = "\n".join(sh.text_frame.text for s in prs.slides for sh in s.shapes if sh.has_text_frame)
        checks = {
            "deck reopens with python-pptx": True,          # reaching here means it parsed
            "title slide + one slide per H2 (3 slides)": len(prs.slides) == 3,
            "title carries the H1": "Platform Update" in text,
            "H2 sections became slides": "Highlights" in text and "Next" in text,
            "list items became bullets": "publish once to a shared volume" in text,
        }
        for name, ok in checks.items():
            print(f"    assert {name}: {'ok' if ok else 'FAILED'}")

        if all(checks.values()):
            print(f"\nRESULT: PASS - branded-pptx ran on serverless and produced a real deck "
                  f"({time.time()-started:.1f}s)")
            return 0
        print("\nRESULT: FAIL - one or more assertions failed")
        return 1
    finally:
        # 5) CLEAN UP (best effort) unless --keep. out_path is None when the run never produced one.
        if not args.keep:
            for p in (in_path, out_path):
                if not p:
                    continue
                try:
                    w.files.delete(p)
                except Exception as e:  # noqa: BLE001 - cleanup is best-effort
                    step(f"cleanup warning for {p}: {e}")
            # Normally absent; present only if the guard quarantined the input. Quiet on purpose.
            for p in (f"{base}/rejected/{stem}.md", f"{base}/rejected/{stem}.md.reason.txt"):
                try:
                    w.files.delete(p)
                except Exception:  # noqa: BLE001 - expected to be missing on the happy path
                    pass
            step("cleaned up test files")


if __name__ == "__main__":
    sys.exit(main())
