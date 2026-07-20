"""
Deterministic markdown -> branded .pptx, pure Python via python-pptx.

MVP-2 in the branded-pptx ladder: it runs on Databricks Free Edition serverless (no Node, no
LibreOffice, no system packages). Lower fidelity than the faithful engine (pptxgenjs + LibreOffice
+ a vision-in-the-loop QA cycle, which MVP-4 restores on paid classic compute): here the layout is
a straight structural map from the markdown, with a light brand skin (accent colour + footer). No
LLM and no vision QA - the deck is a pure function of the markdown, which is exactly what makes it
testable and serverless-safe.

Mapping:
  # H1              -> the title slide (first paragraph after it becomes the subtitle)
  ## H2             -> a new content slide titled with the heading
  - / * bullets     -> bullet points on the current slide
  plain paragraphs  -> bullet points on the current slide (as body text)

Usage:
    python scripts/build_pptx.py <input.md> [-o out.pptx] [--brand coral] [--title "Deck title"]
"""
import argparse
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

# A brand is just an accent colour + a footer label. Deliberately small - MVP-2 is a skin, not a
# full brand system (that is the faithful engine's job). Add brands here without touching the code.
BRANDS = {
    "default": {"accent": "1F4E79", "footer": ""},
    "coral":   {"accent": "E2725B", "footer": "coral"},
    "magma":   {"accent": "B22222", "footer": "Magma"},
}


def parse_markdown(md: str):
    """Split markdown into (title, subtitle, [(section_title, [bullets]), ...]). Deterministic."""
    lines = md.replace("\r\n", "\n").split("\n")
    title, subtitle = "", ""
    sections = []            # list of [section_title, bullets]
    current = None

    for raw in lines:
        line = raw.rstrip()
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith("# ") and not title:
            title = stripped[2:].strip()
            continue
        heading = re.match(r"^#{2,}\s*(.*)$", stripped)   # H2 or deeper starts a section
        if heading:
            current = [heading.group(1).strip(), []]       # title may be empty; build defaults it
            sections.append(current)
            continue
        # a bullet or a plain line
        text = stripped
        m = re.match(r"^[-*]\s+(.*)$", stripped)
        if m:
            text = m.group(1).strip()
        if current is None:
            # content before the first ## : the first such paragraph is the subtitle
            if title and not subtitle and not stripped.startswith("#"):
                subtitle = text
            continue
        current[1].append(text)

    return title, subtitle, sections


def _accent(brand: dict) -> RGBColor:
    return RGBColor.from_string(brand["accent"])


def _set_title(shape, text: str, accent: RGBColor):
    """Set a title placeholder's text and accent-colour it. Guards the empty-title case: an
    empty string produces no run, so colouring runs[0] would IndexError."""
    shape.text = text
    runs = shape.text_frame.paragraphs[0].runs
    if runs:
        runs[0].font.color.rgb = accent


def build_deck(md: str, out_path, brand: str = "default", title_override: str = "") -> str:
    """Build a branded .pptx from markdown text and write it to out_path. Returns out_path.

    out_path is a filesystem path OR any binary file-like object (python-pptx accepts both).
    A caller whose storage cannot seek should pass a buffer and write the bytes itself.
    """
    brand_cfg = BRANDS.get(brand, BRANDS["default"])
    accent = _accent(brand_cfg)
    title, subtitle, sections = parse_markdown(md)
    if title_override:
        title = title_override
    if not title:
        title = "Untitled deck"

    prs = Presentation()

    # Title slide.
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _set_title(slide.shapes.title, title, accent)
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle

    # One content slide per ## section.
    for section_title, bullets in sections:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        _set_title(s.shapes.title, section_title or "Section", accent)
        body = s.placeholders[1].text_frame
        body.clear()
        for i, b in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = b
        _add_footer(s, brand_cfg["footer"], accent)

    prs.save(out_path)
    return out_path


def _add_footer(slide, label: str, accent: RGBColor):
    """A small brand footer in the slide's bottom-left, if the brand defines one."""
    if not label:
        return
    from pptx.util import Inches
    box = slide.shapes.add_textbox(Inches(0.3), Inches(7.0), Inches(3), Inches(0.3))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = label
    run.font.size = Pt(9)
    run.font.color.rgb = accent


def main():
    ap = argparse.ArgumentParser(description="Build a branded .pptx from a markdown file.")
    ap.add_argument("input", help="Path to a markdown file")
    ap.add_argument("-o", "--out", default=None, help="Output .pptx path (default: <input>.pptx)")
    ap.add_argument("--brand", default="default", choices=sorted(BRANDS), help="Brand skin")
    ap.add_argument("--title", default="", help="Override the deck title")
    args = ap.parse_args()

    with open(args.input, "r", encoding="utf-8") as f:
        md = f.read()
    out = args.out or (args.input.rsplit(".", 1)[0] + ".pptx")
    path = build_deck(md, out, brand=args.brand, title_override=args.title)
    print(f"WROTE {path}")


if __name__ == "__main__":
    main()
