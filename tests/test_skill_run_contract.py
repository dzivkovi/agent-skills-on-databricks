"""Contract test for the uniform skill entrypoint (issue #16): every skills/<name>/scripts/run.py
must expose run(ctx) -> str (the path it wrote), regardless of whether the skill is an "analyze"
skill (document-insights, readability - deterministic metrics + one LLM call) or a "builder" skill
(branded-pptx - no LLM call, writes a python-pptx deck).

Creds-free by design: ctx["llm"] is a stub that never touches a real serving endpoint, so this
runs anywhere with no Databricks workspace. Loads each skill's scripts/run.py by file path, exactly
as src/run_skill.py does, so a skill published alone on a UC volume (no repo on the Python path,
per #6) is exercised the same way it will be in production.
"""
import importlib.util
import logging
from pathlib import Path

import pytest
from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
SKILLS = sorted(
    p for p in (ROOT / "skills").iterdir()
    if p.is_dir() and (p / "SKILL.md").exists()
)
ANALYZE_SKILLS = [p for p in SKILLS if (p / "scripts" / "analyze.py").exists()]
BUILDER_SKILLS = [p for p in SKILLS if not (p / "scripts" / "analyze.py").exists()]

# A few sentences plus an H1 and an H2 with two bullets, so both shapes get something real to
# work with: analyze skills get sentences/words to count, branded-pptx gets a title + a section.
SAMPLE_TEXT = (
    "# Sample Doc\n\n"
    "This is the first sentence of the sample document. Here is a second sentence for good "
    "measure. And a third one so the deterministic metrics have something real to count.\n\n"
    "## Points\n\n"
    "- First point about the release.\n"
    "- Second point about the rollout.\n"
)
TODAY = "2026-07-16"
IN_PATH = "/in/sample-doc.md"


class _StubLLM:
    """Records every call and returns a fixed, recognizable reading - never touches a real
    endpoint, so the whole suite runs creds-free."""

    def __init__(self):
        self.calls = []

    def __call__(self, messages, max_tokens=1500):
        self.calls.append((messages, max_tokens))
        return "STUB-READING sentiment positive."


def _load_run(skill_dir: Path):
    """Import a skill's run(ctx) by file path - the same mechanism src/run_skill.py uses, and
    the only one that also works for a skill published alone on a UC volume (no repo on the
    Python path)."""
    spec = importlib.util.spec_from_file_location(f"{skill_dir.name}_run", skill_dir / "scripts" / "run.py")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module.run


def _make_ctx(skill_dir: Path, tmp_path: Path, llm: _StubLLM) -> dict:
    """A frozen ctx exactly matching the contract every skill's run.py receives."""
    out_dir = str(tmp_path)
    skill_name = skill_dir.name
    stem = "sample-doc"
    out_base = f"{out_dir}/{stem}-{skill_name}-{TODAY}"   # no extension - that's the skill's call
    return {
        "text": SAMPLE_TEXT,
        "in_path": IN_PATH,
        "out_dir": out_dir,
        "out_base": out_base,
        "skill_dir": str(skill_dir),
        "skill_md": (skill_dir / "SKILL.md").read_text(encoding="utf-8"),
        "skill_name": skill_name,
        "model": "test-endpoint",
        "llm": llm,
        "log": logging.getLogger(f"test.{skill_name}"),
        "today": TODAY,
    }


@pytest.mark.parametrize("skill_dir", SKILLS, ids=lambda p: p.name)
def test_run_writes_a_namespaced_file_inside_out_dir(skill_dir, tmp_path):
    llm = _StubLLM()
    ctx = _make_ctx(skill_dir, tmp_path, llm)
    run = _load_run(skill_dir)

    out_path = run(ctx)

    assert isinstance(out_path, str), "run(ctx) must return the path it wrote, as a str"
    written = Path(out_path)
    assert written.is_file() and written.stat().st_size > 0
    assert written.resolve().parent == tmp_path.resolve(), "must write inside ctx['out_dir']"
    expected_prefix = f"sample-doc-{skill_dir.name}-{TODAY}"
    assert written.name.startswith(expected_prefix), \
        f"expected {written.name} to start with {expected_prefix!r} (the collision-proof namespacing)"


@pytest.mark.parametrize("skill_dir", ANALYZE_SKILLS, ids=lambda p: p.name)
def test_analyze_skill_run_writes_a_labelled_markdown_report(skill_dir, tmp_path):
    llm = _StubLLM()
    ctx = _make_ctx(skill_dir, tmp_path, llm)
    run = _load_run(skill_dir)

    out_path = run(ctx)

    assert out_path.endswith(".md")
    content = Path(out_path).read_text(encoding="utf-8")
    assert "## Metrics (computed by code - exact)" in content
    assert "## Reading (interpreted by the LLM)" in content
    assert "STUB-READING" in content
    assert "sample-doc.md" in content, "the report must cite its source file for provenance"
    assert len(llm.calls) == 1, "an analyze skill must call the LLM exactly once"


@pytest.mark.parametrize("skill_dir", BUILDER_SKILLS, ids=lambda p: p.name)
def test_builder_skill_run_writes_an_openable_deck_with_no_llm_call(skill_dir, tmp_path):
    llm = _StubLLM()
    ctx = _make_ctx(skill_dir, tmp_path, llm)
    run = _load_run(skill_dir)

    out_path = run(ctx)

    assert out_path.endswith(".pptx")
    assert llm.calls == [], "a builder skill is deterministic - it must never call the LLM"
    prs = Presentation(out_path)
    assert len(prs.slides) >= 2
    title_slide = list(prs.slides)[0]
    assert "Sample Doc" in title_slide.shapes.title.text


def test_all_skills_produce_pairwise_distinct_output_paths(tmp_path):
    # The namespacing invariant, now proven at the run() level rather than at output_base() alone:
    # the same fake input through every skill must never collide on disk.
    paths = []
    for skill_dir in SKILLS:
        llm = _StubLLM()
        ctx = _make_ctx(skill_dir, tmp_path, llm)
        run = _load_run(skill_dir)
        paths.append(run(ctx))
    assert len(paths) == len(set(paths)), f"expected pairwise-distinct output paths, got {paths}"
