"""Contract test for the harness's own adapters (issue #21): src/adapters.py exposes ADAPTERS =
{"report": report, "deck": deck}, and src/run_skill.py's load_adapter(name) resolves one by name.

This replaces tests/test_skill_run_contract.py, which loaded each skill's own scripts/run.py by
file path - that file no longer exists. Imported skills are now READ-ONLY (SKILL.md plus their own
scripts/, nothing harness-shaped); the Databricks-specific plumbing moved here, into the adapters,
so this is the test that owns proving the adapter contract instead of the skill's.

Creds-free by design: ctx["llm"] is a stub that never touches a real serving endpoint, so this
runs anywhere with no Databricks workspace.
"""
import logging
from pathlib import Path

import pytest
from pptx import Presentation

import adapters
import run_skill
from adapters import ADAPTERS

ROOT = Path(__file__).resolve().parents[1]
SKILLS = sorted(
    p for p in (ROOT / "skills").iterdir()
    if p.is_dir() and (p / "SKILL.md").exists()
)
ANALYZE_SKILLS = [p for p in SKILLS if (p / "scripts" / "analyze.py").exists()]
BUILDER_SKILLS = [p for p in SKILLS if (p / "scripts" / "build_pptx.py").exists()]

# A few sentences plus an H1 and an H2 with two bullets, so both adapter shapes get something real
# to work with: report gets sentences/words to count, deck gets a title + a section.
SAMPLE_TEXT = (
    "# Sample Doc\n\n"
    "This is the first sentence of the sample document. Here is a second sentence for good "
    "measure. And a third one so the deterministic metrics have something real to count.\n\n"
    "## Points\n\n"
    "- First point about the release.\n"
    "- Second point about the rollout.\n"
)
TODAY = "2026-07-16"
IN_PATH = "/in/sample-doc.md"


class _StubLLM:
    """Records every call and returns a fixed, recognizable reading - never touches a real
    endpoint, so the whole suite runs creds-free."""

    def __init__(self):
        self.calls = []

    def __call__(self, messages, max_tokens=1500):
        self.calls.append((messages, max_tokens))
        return "STUB-READING sentiment positive."


def _make_ctx(skill_dir: Path, tmp_path: Path, llm: _StubLLM) -> dict:
    """A frozen ctx exactly matching the contract every adapter receives - unchanged by #21."""
    out_dir = str(tmp_path)
    skill_name = skill_dir.name
    stem = "sample-doc"
    out_base = f"{out_dir}/{stem}-{skill_name}-{TODAY}"   # no extension - the adapter appends its own
    return {
        "text": SAMPLE_TEXT,
        "in_path": IN_PATH,
        "out_dir": out_dir,
        "out_base": out_base,
        "skill_dir": str(skill_dir),
        "skill_md": (skill_dir / "SKILL.md").read_text(encoding="utf-8"),
        "skill_name": skill_name,
        "model": "test-endpoint",
        "llm": llm,
        "log": logging.getLogger(f"test.{skill_name}"),
        "today": TODAY,
    }


@pytest.mark.parametrize("skill_dir", ANALYZE_SKILLS, ids=lambda p: p.name)
def test_report_adapter_writes_a_labelled_markdown_report(skill_dir, tmp_path):
    llm = _StubLLM()
    ctx = _make_ctx(skill_dir, tmp_path, llm)

    out_path = adapters.report(ctx)

    assert isinstance(out_path, str) and out_path.endswith(".md")
    written = Path(out_path)
    assert written.resolve().parent == tmp_path.resolve(), "must write inside ctx['out_dir']"
    assert written.name == f"sample-doc-{skill_dir.name}-{TODAY}.md", \
        "expected the collision-proof <stem>-<skill>-<today>.md namespacing"
    content = written.read_text(encoding="utf-8")
    assert "## Metrics (computed by code - exact)" in content
    assert "## Reading (interpreted by the LLM)" in content
    assert "STUB-READING" in content
    assert "sample-doc.md" in content, "the report must cite its source file for provenance"
    assert len(llm.calls) == 1, "the report adapter must call the LLM exactly once"


def test_deck_adapter_writes_an_openable_deck_with_no_llm_call(tmp_path):
    skill_dir = next(p for p in BUILDER_SKILLS if p.name == "branded-pptx")
    llm = _StubLLM()
    ctx = _make_ctx(skill_dir, tmp_path, llm)

    out_path = adapters.deck(ctx)

    assert out_path.endswith(".pptx")
    assert llm.calls == [], "the deck adapter is deterministic - it must never call the LLM"
    prs = Presentation(out_path)
    assert len(prs.slides) >= 2
    title_slide = list(prs.slides)[0]
    assert "Sample Doc" in title_slide.shapes.title.text


def test_report_over_both_analyze_skills_yields_distinct_paths(tmp_path):
    # The namespacing invariant: the same fake input through both analyze skills must never
    # collide on disk. Proven at the adapter level now that skills carry no run.py of their own.
    paths = [adapters.report(_make_ctx(skill_dir, tmp_path, _StubLLM())) for skill_dir in ANALYZE_SKILLS]
    assert len(paths) == len(set(paths)), f"expected pairwise-distinct output paths, got {paths}"


def test_load_adapter_resolves_known_names_to_callables():
    assert run_skill.load_adapter("report") is ADAPTERS["report"]
    assert run_skill.load_adapter("deck") is ADAPTERS["deck"]
    assert callable(run_skill.load_adapter("report"))
    assert callable(run_skill.load_adapter("deck"))


def test_load_adapter_rejects_unknown_name():
    with pytest.raises(SystemExit):
        run_skill.load_adapter("nope")


def test_wrong_adapter_for_a_skill_fails_loudly_not_silently(tmp_path):
    # The new failure mode --adapter introduces: an adapter loads a skill's script BY NAME
    # (scripts/analyze.py or scripts/build_pptx.py), so pointing it at a skill that doesn't ship
    # that script must raise, never silently produce nothing or the wrong shape.
    document_insights = next(p for p in ANALYZE_SKILLS if p.name == "document-insights")
    branded_pptx = next(p for p in BUILDER_SKILLS if p.name == "branded-pptx")

    with pytest.raises(SystemExit):
        adapters.deck(_make_ctx(document_insights, tmp_path, _StubLLM()))

    with pytest.raises(SystemExit):
        adapters.report(_make_ctx(branded_pptx, tmp_path, _StubLLM()))
