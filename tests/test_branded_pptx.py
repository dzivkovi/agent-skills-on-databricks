"""Integration test for the branded-pptx skill (#2): build a REAL .pptx from markdown and prove
it by reopening it with python-pptx and asserting the slides, titles, bullets, and brand accent.

Self-contained: loads the skill's build script by file path (no conftest), so it never collides
with another branch's test setup. This is the full workflow for this skill - markdown in, a valid
.pptx out - runnable with no Databricks workspace.
"""
import importlib.util
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parents[1]
SKILL = ROOT / "skills" / "branded-pptx"
SAMPLE = SKILL / "samples" / "quarterly-review.md"


def _load_builder():
    spec = importlib.util.spec_from_file_location(
        "build_pptx", SKILL / "scripts" / "build_pptx.py")
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


MD = _load_builder()  # module handle reused across tests


def _all_text(prs) -> str:
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
    return "\n".join(out)


def test_build_produces_a_valid_openable_pptx(tmp_path):
    out = str(tmp_path / "deck.pptx")
    md = SAMPLE.read_text(encoding="utf-8")
    path = MD.build_deck(md, out, brand="coral")

    # It is a real file, and python-pptx can reopen it (so PowerPoint/Slides can too).
    assert Path(path).is_file() and Path(path).stat().st_size > 0
    prs = Presentation(path)

    # Title slide + one slide per ## section (3 sections in the sample) = 4 slides.
    assert len(prs.slides) == 4

    text = _all_text(prs)
    assert "Q3 Platform Review" in text              # H1 -> title slide
    assert "Highlights" in text and "Reliability" in text and "Next Quarter" in text  # H2 titles
    assert "Migrated ingestion to serverless" in text[:len(text)] and "40 percent" in text  # a bullet
    assert "Zero Sev-1 incidents this quarter" in text


def test_title_slide_carries_h1_and_subtitle(tmp_path):
    out = str(tmp_path / "d.pptx")
    MD.build_deck("# My Title\n\nA one line subtitle.\n\n## Section A\n\n- point one\n", out)
    prs = Presentation(out)
    title_slide = list(prs.slides)[0]
    assert title_slide.shapes.title.text == "My Title"
    # subtitle placeholder carries the first paragraph after the H1
    assert any(s.has_text_frame and "one line subtitle" in s.text_frame.text
               for s in title_slide.shapes)


def test_brand_accent_colour_applied_to_titles(tmp_path):
    out = str(tmp_path / "d.pptx")
    MD.build_deck("# T\n\n## S\n\n- b\n", out, brand="coral")
    prs = Presentation(out)
    title_run = list(prs.slides)[0].shapes.title.text_frame.paragraphs[0].runs[0]
    assert title_run.font.color.rgb == RGBColor.from_string(MD.BRANDS["coral"]["accent"])


def test_bullets_map_from_list_items(tmp_path):
    out = str(tmp_path / "d.pptx")
    MD.build_deck("# T\n\n## S\n\n- alpha\n- beta\n- gamma\n", out)
    prs = Presentation(out)
    section = list(prs.slides)[1]
    body = [s for s in section.shapes if s.has_text_frame and s != section.shapes.title][0]
    lines = [p.text for p in body.text_frame.paragraphs]
    assert lines == ["alpha", "beta", "gamma"]


def test_untitled_when_no_h1(tmp_path):
    out = str(tmp_path / "d.pptx")
    MD.build_deck("## Only a section\n\n- x\n", out)
    prs = Presentation(out)
    assert list(prs.slides)[0].shapes.title.text == "Untitled deck"


def test_empty_h2_heading_is_kept_as_a_section_not_dropped(tmp_path):
    # Regression: '## ' (empty title) must still create a section (was silently dropped, and its
    # bullets leaked to the previous context) and must not crash on the empty title run.
    out = str(tmp_path / "d.pptx")
    MD.build_deck("# T\n\n## \n\n- a bullet under the blank heading\n", out)
    prs = Presentation(out)
    assert len(prs.slides) == 2                       # title + the (blank-titled) section
    section = list(prs.slides)[1]
    assert section.shapes.title.text == "Section"     # empty title defaulted, no crash
    assert "a bullet under the blank heading" in _all_text(prs)


def test_deeper_headings_start_sections_too(tmp_path):
    out = str(tmp_path / "d.pptx")
    MD.build_deck("# T\n\n## Two\n\n- a\n\n### Three\n\n- b\n", out)
    prs = Presentation(out)
    assert len(prs.slides) == 3                        # title + H2 + H3
    assert "Three" in _all_text(prs) and "b" in _all_text(prs)
